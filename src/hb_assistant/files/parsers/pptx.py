"""PPTXParser: bounded slide text + notes extraction (python-pptx). No media/ole extraction."""

from __future__ import annotations

from pathlib import Path
from typing import Any, Dict

try:
    from pptx import Presentation  # type: ignore
    from pptx.util import Inches, Pt  # type: ignore
except ImportError:
    Presentation = None  # type: ignore


class PPTXParser:
    MAX_SLIDES = 20

    def parse(self, path: Path, max_chars: int = 8000) -> Dict[str, Any]:
        if Presentation is None:
            raise RuntimeError("python-pptx not installed")
        try:
            prs = Presentation(str(path))
            parts: list[str] = []
            total = 0
            slide_count = 0
            for slide in prs.slides:
                slide_count += 1
                if slide_count > self.MAX_SLIDES:
                    break
                parts.append(f"--- slide {slide_count} ---")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        t = shape.text.strip()
                        if t:
                            parts.append(t)
                            total += len(t)
                            if total > max_chars:
                                break
                    if hasattr(shape, "text_frame"):
                        for para in shape.text_frame.paragraphs:
                            t = para.text.strip()
                            if t:
                                parts.append(t)
                                total += len(t)
                                if total > max_chars:
                                    break
                # notes
                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text or ""
                    if notes.strip():
                        parts.append(f"[notes] {notes.strip()}")
                        total += len(notes)
                if total > max_chars:
                    break
            excerpt = "\n".join(parts)[:max_chars]
            return {
                "text_excerpt": excerpt,
                "char_count": len(excerpt),
                "slide_count": len(prs.slides),
                "slides_sampled": min(slide_count, self.MAX_SLIDES),
            }
        except Exception as e:
            msg = str(e)[:200]
            fc = "parser_error"
            return {"text_excerpt": "", "char_count": 0, "error": msg, "failure_code": fc}
