"""N8C-24 — generated-output file writers.

Real generation for every allowed type (docx/xlsx/pptx/pdf via python-docx/openpyxl/python-pptx/reportlab;
html/csv/json/txt/md as validated text; zip validated then written). No format is faked. Every write is
atomic: render to a temp file in the destination directory, fsync, then ``os.replace`` onto the final path.
The caller (workspace repo) is responsible for approval/idempotency/receipts; this module only renders bytes.
"""

from __future__ import annotations

import base64
import csv
import hashlib
import io
import json
import os
import tempfile
from pathlib import Path
from typing import Any

from .client_output_zip import validate_zip_payload
from .config import NasMcpConfig


class OutputWriteError(Exception):
    """Generated-output rendering or write failed."""


def sha256_hex(data: bytes) -> str:
    return hashlib.sha256(data).hexdigest()


def _atomic_write_bytes(target: Path, data: bytes) -> None:
    target.parent.mkdir(parents=True, exist_ok=True)
    fd, tmp = tempfile.mkstemp(dir=str(target.parent), suffix=".tmp")
    try:
        with os.fdopen(fd, "wb") as fh:
            fh.write(data)
            fh.flush()
            os.fsync(fh.fileno())
        os.replace(tmp, str(target))
    finally:
        if os.path.exists(tmp):
            os.unlink(tmp)


def _decode_base64(content: str) -> bytes:
    try:
        return base64.b64decode(content, validate=True)
    except (ValueError, TypeError) as exc:
        raise OutputWriteError(f"invalid base64 payload: {exc}") from exc


# --- per-type renderers (bytes, no IO) ---

def _render_text(content: str) -> bytes:
    return str(content).encode("utf-8")


def _render_json(content: str) -> bytes:
    try:
        json.loads(content)
    except (ValueError, TypeError) as exc:
        raise OutputWriteError(f"invalid json content: {exc}") from exc
    return str(content).encode("utf-8")


def _render_csv(content: str) -> bytes:
    # Accept CSV text as-is (already delimited); just ensure it parses.
    list(csv.reader(io.StringIO(str(content))))
    return str(content).encode("utf-8")


def _render_docx(content: str) -> bytes:
    from docx import Document  # noqa: PLC0415

    doc = Document()
    for line in str(content).splitlines() or [""]:
        if line.startswith("# "):
            doc.add_heading(line[2:], level=1)
        elif line.startswith("## "):
            doc.add_heading(line[3:], level=2)
        else:
            doc.add_paragraph(line)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _render_xlsx(content: str) -> bytes:
    from openpyxl import Workbook  # noqa: PLC0415

    wb = Workbook()
    ws = wb.active
    for row_idx, row in enumerate(csv.reader(io.StringIO(str(content))), start=1):
        for col_idx, value in enumerate(row, start=1):
            ws.cell(row=row_idx, column=col_idx, value=value)
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _render_pptx(content: str) -> bytes:
    from pptx import Presentation  # noqa: PLC0415

    prs = Presentation()
    # Each blank line separates slides; first line of a block is the title, the rest is the body.
    blocks = [b.strip() for b in str(content).split("\n\n") if b.strip()] or ["Slide"]
    layout = prs.slide_layouts[1]
    for block in blocks:
        lines = block.splitlines()
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = lines[0][:250]
        if len(lines) > 1:
            slide.placeholders[1].text = "\n".join(lines[1:])
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _render_html(content: str) -> bytes:
    return str(content).encode("utf-8")


def _render_pdf(content: str) -> bytes:
    from reportlab.lib.pagesizes import LETTER  # noqa: PLC0415
    from reportlab.lib.styles import getSampleStyleSheet  # noqa: PLC0415
    from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer  # noqa: PLC0415

    styles = getSampleStyleSheet()
    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=LETTER)
    flow: list[Any] = []
    for line in str(content).splitlines():
        if not line.strip():
            flow.append(Spacer(1, 8))
            continue
        style = styles["Heading1"] if line.startswith("# ") else (
            styles["Heading2"] if line.startswith("## ") else styles["BodyText"])
        text = line.lstrip("# ").replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
        flow.append(Paragraph(text or "&nbsp;", style))
    if not flow:
        flow.append(Paragraph("&nbsp;", styles["BodyText"]))
    doc.build(flow)
    return buf.getvalue()


def _render_zip(config: NasMcpConfig, content: str) -> tuple[bytes, dict[str, Any]]:
    data = _decode_base64(content)
    summary = validate_zip_payload(config, data)  # raises on any unsafe member
    return data, summary


# content_mode → (extension, renderer)
def render_output_bytes(
    *, config: NasMcpConfig, file_type: str, content_mode: str, content: str,
) -> tuple[bytes, dict[str, Any]]:
    """Render a generated output to bytes. Returns (bytes, extra) where extra may carry zip_validation."""
    ext = str(file_type).lower()
    extra: dict[str, Any] = {}
    if ext in ("txt", "md"):
        data = _render_text(content)
    elif ext == "html":
        data = _render_html(content)
    elif ext == "json":
        data = _render_json(content)
    elif ext == "csv":
        data = _render_csv(content)
    elif ext == "docx":
        data = _decode_base64(content) if content_mode == "base64_binary" else _render_docx(content)
    elif ext == "xlsx":
        data = _decode_base64(content) if content_mode == "base64_binary" else _render_xlsx(content)
    elif ext == "pptx":
        data = _decode_base64(content) if content_mode == "base64_binary" else _render_pptx(content)
    elif ext == "pdf":
        data = _decode_base64(content) if content_mode == "base64_binary" else _render_pdf(content)
    elif ext == "zip":
        data, extra["zip_validation"] = _render_zip(config, content)
    else:
        raise OutputWriteError(f"no renderer for file_type: {ext}")
    if len(data) > config.max_client_output_file_bytes:
        raise OutputWriteError("rendered output exceeds max_client_output_file_bytes")
    return data, extra


def write_output_bytes(target: Path, data: bytes) -> dict[str, Any]:
    """Atomically write already-rendered bytes to the resolved target path."""
    _atomic_write_bytes(target, data)
    return {"bytes_written": len(data), "sha256": sha256_hex(data)}


def assemble_zip_from_files(config: NasMcpConfig, files: list[tuple[str, Path]]) -> tuple[bytes, dict[str, Any]]:
    """Build a ZIP from already-committed output files (arcname, source path). Sanitizes member names,
    then re-validates the assembled bytes through the same ZIP validator."""
    import zipfile  # noqa: PLC0415

    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", compression=zipfile.ZIP_DEFLATED) as zf:
        for arcname, src in files:
            safe = str(arcname).replace("\\", "/").lstrip("/").replace("..", "")
            if not src.is_file():
                raise OutputWriteError(f"zip source not a file: {src}")
            zf.writestr(safe, src.read_bytes())
    data = buf.getvalue()
    summary = validate_zip_payload(config, data)
    return data, summary
